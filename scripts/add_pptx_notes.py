"""
Adds speaker notes to all 14 slides of BoxNiJuanPH_OralPresentation_June21.pptx.
Preserves all existing slide content (text, shapes, design) — only adds to notes pane.
"""

from pptx import Presentation
from pptx.util import Pt
import copy

path = r"C:\Users\I331629\Documents\BoxNiJuanPH-guide\BoxNiJuanPH_OralPresentation_June21.pptx"
prs = Presentation(path)

# Speaker notes per slide (0-indexed)
# Written as cue cards — key sentences the presenter can glance at.
NOTES = {
    0: """[SLIDE 1 — TITLE]
RIZZA opens:
Good morning/afternoon po, Professor Carabaca! Kaming limang miyembro ng BSITOUMN 2-1 ay nagtatanghal ngayon ng aming COMP 047 project — BoxNiJuanPH: An Interactive Wellness Subscription Box Platform.

Ang aming live site ay accessible ngayon sa boxnijuanph.vercel.app. Sisimulan namin ngayon ang aming presentation.
""",

    1: """[SLIDE 2 — OUTLINE]
RIZZA continues:
Ang aming presentation ngayon ay nahahati sa limang bahagi:
Una — ang problem background at introduction.
Pangalawa — ang aming anim na SMART objectives.
Pangatlo — ang comparison namin sa related systems.
Pang-apat — ang technical specifications ng aming platform.
At panghuli — ang live prototype demonstration, na ipapakita ni Lacel.

Note: Si Ivy ay hindi makakasama ngayon dahil sa power interruption sa kanilang lugar — ang kanyang bahagi ay ibinabahagi namin sa aming group.
""",

    2: """[SLIDE 3 — INTRODUCTION & BACKGROUND]
CARLA:
Ito po ang problema na sinolve namin. Most subscription box services — pre-curated lahat ng laman. Hindi mo pinipili kung ano ang nasa box mo. Kaya naman mataas ang cancellation rate.

Dagdag pa, international platforms like Cratejoy or FabFitFun — walang Philippine Peso pricing, walang GCash or Maya, at walang local Filipino brands.

Ayon sa Grand View Research 2023, ang global subscription box market ay nagkakahalaga ng $32 billion in 2022, at lalaki pa ito hanggang $105 billion by 2032. At ayon sa Epsilon 2022, 80% ng consumers ay mas likely bumili kung personalized ang experience.

Kaya ang BoxNiJuanPH — binaliktad namin ang modelo. Ikaw mismo ang pumipili ng bawat item sa iyong monthly wellness box.
""",

    3: """[SLIDE 4 — OBJECTIVES]
CARLA:
Ang aming anim na SMART objectives:

O1 — Users can build and order a personalized box mula sa 20 products across 4 categories, by June 14, 2026.
O2 — Apat na subscription tiers: Basic sa ₱399, Standard sa ₱599, Premium sa ₱899, at Custom sa ₱1,299 — lahat competitive within Philippine market pricing.
O3 — Easy-to-navigate UX para sa Filipino consumers na 18 to 35 years old.
O4 — CSR integration — Filipino Brand tags, Eco-Friendly badges, at CSR impact summary after every order, aligned sa UN SDG 12.
O5 — Credible market positioning gamit ang 4 or more na cited research references.
O6 — Complete submission — ang buong working platform, documentation, at presentation, submitted by June 14, 2026. Natupad po namin ito.
""",

    4: """[SLIDE 5 — RELATED SYSTEMS]
JULIO:
Pinaghalintulad namin ang BoxNiJuanPH sa existing platforms.

Cratejoy at FabFitFun — US-based. Walang PHP pricing, walang local brands, at limited lang ang customization.
Good Box PH — nasa Pilipinas sila pero fully pre-curated. Walang kahit anong customization at walang subscription management dashboard.
Boozy.ph at BeautyMNL — hindi wellness subscription box ang pangunahing focus nila.

BoxNiJuanPH ang tanging platform sa Philippine market na nag-ooffer ng: full item-level customization sa bawat plan tier, Filipino brand curation, PHP pricing, subscription management tools, at built-in CSR integration.
""",

    5: """[SLIDE 6 — TECH SPECS: ARCHITECTURE]
JULIO:
Ang aming tech stack: Next.js App Router gamit ang TypeScript at Tailwind CSS v4. Deployed sa Vercel Hobby Plan — libre.

Ang lahat ng data ay stored sa user's device lang — walang external server, walang database. RA 10173 compliant ito dahil walang personal data na nai-send sa labas ng device.

Total development cost: ₱0. Live ngayon sa boxnijuanph.vercel.app.
""",

    6: """[SLIDE 7 — TECH SPECS: 11 PAGES]
JULIO:
Ang aming platform ay may 11 fully responsive pages — mula sa homepage hanggang privacy policy.

Ang pinakamahalagang pages ay: ang Box Builder — dito pumipili ang user ng products — at ang My Box Dashboard — dito nila mina-manage ang kanilang subscription.

Ang Custom plan sa /plans ay nag-aalok ng up to 12 items per box — pinaka-flexible na tier.
""",

    7: """[SLIDE 8 — USE CASES]
JULIO:
Apat na use cases ang aming na-document.

UC1 — Registered user, Standard plan — 5 items, checkout gamit ang pre-filled form, order number, at active subscription sa My Box.
UC2 — Guest checkout — may amber disclaimer ang limitations, manual form, tapos nailagay ang order.
UC3 — Returning user — nag-edit ng profile, nag-set ng default GCash — auto-fill na sa susunod na checkout.
UC4 — Cancel subscription — dalawang confirm steps lang, walang lock-in, walang phone call.
""",

    8: """[SLIDE 9 — CSR & SUSTAINABILITY]
CARLA:
Hindi lang feature ang CSR namin — ito ay embedded sa buong experience.

Una — every locally sourced product ay may Filipino Brand badge. Nakikita ito habang nag-browse at sa bawat order summary.
Pangalawa — dahil ikaw mismo ang pumipili ng items, walang unwanted products, walang basura. Direkta itong application ng UN SDG 12 — Responsible Consumption and Production.
Pangatlo — sa Box Summary, Order Confirmation, at My Box — nakikita ng user kung ilang local at eco-friendly items ang nasa box nila.

Lahat ng CSR stats ay real-time — hindi lang decorative.
""",

    9: """[SLIDE 10 — DEMO INTRO]
LACEL:
Good afternoon po ulit, Professor. Ngayon ay ipapakita ko ang live platform — ang BoxNiJuanPH — sa boxnijuanph.vercel.app.

Ang demo namin ay sumasaklaw ng 8 steps: Homepage, Plan Selection, Box Builder, Box Summary, Login or Guest checkout, Checkout, Order Confirmation, at My Box Dashboard.

Magsisimula na tayo.
""",

    10: """[SLIDE 11 — DEMO: HOMEPAGE & PLANS]
LACEL:
Ito ang Homepage. Makikita ninyo ang hero section, ang How It Works steps, featured products, at ang Plans Preview section.

Dito rin namin idinagdag ang Wellness Quiz — tatlong tanong lang, at irerecommend nito ang tamang plan para sa user.

Ngayon ay pupunta tayo sa /plans. Makikita ninyo ang apat na tiers — Basic, Standard, Premium, at Custom — pati na ang savings badges at loyalty perks strip sa ibaba.
""",

    11: """[SLIDE 12 — DEMO: BUILDER & SUMMARY]
LACEL:
Ito ang Box Builder. Magse-select tayo ng plan at pipili ng items.

Makikita ninyo ang: category at brand filters, ang variant modal kapag nag-click ng product, at ang progress bar na nagse-show kung ilan na ang napili.

Ang Custom plan — up to 12 items ang pwedeng ilagay. Merong ⓘ button sa bawat card para makita ang full product details nang hindi mina-select ang item.

Pagkatapos mapuno ang box, pupunta tayo sa /summary para ma-review ang lahat ng items at ang CSR impact count.
""",

    12: """[SLIDE 13 — DEMO: CHECKOUT & CONFIRMATION]
LACEL:
Sa /checkout — makikita ang delivery form na may validation: PH phone number format, address fields, at lahat ng required fields ay nire-require bago pwedeng i-submit.

Pagkatapos mag-confirm, pupunta tayo sa /confirmation — may order number sa BNJ-XXXXXX format, full item list, at ang amber refund notice card — nagpapaalam sa user na may 7-day window sila para mag-request ng refund o replacement.

Mula dito ay mapupunta tayo sa /my-box para makita ang active subscription.
""",

    13: """[SLIDE 14 — CLOSING / Q&A]
RIZZA closes:
Iyon po ang aming BoxNiJuanPH — a fully functional, live, Filipino-first wellness subscription box platform.

Natupad namin ang lahat ng anim na SMART objectives, na-deploy namin ito sa zero cost, at accessible ito ngayon sa boxnijuanph.vercel.app.

Maraming salamat po, Professor Carabaca! Handa na kaming sumagot sa inyong mga tanong. Maaari na rin ninyong subukan ang platform — five minutes lang po para makumpleto ang isang order.
""",
}


def set_slide_notes(slide, text):
    """Set the notes text for a slide, creating the notes slide if needed."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear existing paragraphs
    from pptx.oxml.ns import qn as pqn
    from lxml import etree
    # Remove all existing <a:p> elements inside <p:txBody>
    txBody = tf._txBody
    for p in txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        p.getparent().remove(p)
    # Add our text as paragraphs
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for line in text.split('\n'):
        p_elem = etree.SubElement(txBody, f'{{{ns}}}p')
        r_elem = etree.SubElement(p_elem, f'{{{ns}}}r')
        rPr = etree.SubElement(r_elem, f'{{{ns}}}rPr', attrib={'lang': 'en-PH', 'dirty': '0'})
        t_elem = etree.SubElement(r_elem, f'{{{ns}}}t')
        t_elem.text = line


for i, slide in enumerate(prs.slides):
    if i in NOTES:
        set_slide_notes(slide, NOTES[i])
        print(f"Slide {i+1}: notes added ({len(NOTES[i])} chars)")

prs.save(path)
print("\nDone:", path)
